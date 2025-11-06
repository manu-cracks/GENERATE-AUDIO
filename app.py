# Importaciones necesarias
from flask import Flask, render_template, request, send_file, jsonify
from dotenv import load_dotenv
from openai import OpenAI
from pathlib import Path
import os
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from werkzeug.utils import secure_filename
import uuid

# Cargar todas las variables del archivo .env
load_dotenv()

# Inicializar Flask
app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16MB max-limit
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['OUTPUT_FOLDER'] = 'outputs'

# Crear carpetas si no existen
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
os.makedirs(app.config['OUTPUT_FOLDER'], exist_ok=True)

ALLOWED_EXTENSIONS = {'pdf', 'docx', 'pptx', 'txt'} 

# --- CONFIGURACIÓN DE LAS CLAVES ---

# 1. Obtener la clave de OpenRouter directamente del entorno
# La usaremos para autenticar el cliente de texto.
openrouter_key = os.getenv("OPENROUTER_API_KEY") 

# 2. La clave de OpenAI (OPENAI_API_KEY) se carga automáticamente por defecto.

# --- CLIENTES API ---

# Cliente 1: Para la generación de TEXTO (usando OpenRouter)
# Especificamos la URL base de OpenRouter y la clave
client_text = OpenAI(
    base_url="https://openrouter.ai/api/v1",
    api_key=openrouter_key,
)

# Cliente 2: Para la generación de AUDIO (usando OpenAI)
# Usa la URL base por defecto de OpenAI y la clave OPENAI_API_KEY
client_audio = OpenAI() 

# --- LÓGICA DE EXTRACCIÓN DE TEXTO (mejorada) ---

def allowed_file(filename):
    """Verifica si el archivo tiene una extensión permitida."""
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extraer_texto_pdf(ruta_archivo):
    """Extrae texto de un archivo PDF."""
    try:
        reader = PdfReader(ruta_archivo)
        texto_completo = ""
        for page in reader.pages:
            texto_completo += page.extract_text() or ""
        return texto_completo
    except Exception as e:
        return f"Error al leer PDF: {e}"

def extraer_texto_word(ruta_archivo):
    """Extrae texto de un archivo Word (.docx)."""
    try:
        doc = Document(ruta_archivo)
        texto_completo = ""
        for paragraph in doc.paragraphs:
            texto_completo += paragraph.text + "\n"
        return texto_completo
    except Exception as e:
        return f"Error al leer Word: {e}"

def extraer_texto_ppt(ruta_archivo):
    """Extrae texto de un archivo PowerPoint (.pptx)."""
    try:
        prs = Presentation(ruta_archivo)
        texto_completo = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texto_completo += shape.text + "\n"
        return texto_completo
    except Exception as e:
        return f"Error al leer PowerPoint: {e}"

def extraer_texto(ruta_archivo):
    """Extrae texto de un archivo según su extensión."""
    try:
        extension = ruta_archivo.lower().split('.')[-1]
        
        if extension == 'pdf':
            return extraer_texto_pdf(ruta_archivo)
        elif extension == 'docx':
            return extraer_texto_word(ruta_archivo)
        elif extension == 'pptx':
            return extraer_texto_ppt(ruta_archivo)
        elif extension == 'txt':
            with open(ruta_archivo, 'r', encoding='utf-8') as f:
                return f.read()
        else:
            return "Formato no soportado."
            
    except Exception as e:
        return f"Error al leer el archivo: {e}"

# --- FUNCIÓN PRINCIPAL DE PROCESAMIENTO ---

def procesar_y_generar_audio(contenido_archivo, nombre_salida):
    
    # 1. --- ETAPA DE TEXTO: Generación con OpenRouter ---
    print("Iniciando generación de texto con OpenRouter...")

    # Define el modelo de OpenRouter que quieres usar
    # MODELO ACTUALIZADO: Claude 3.5 Sonnet (excelente para narrativa)
    MODELO_TEXTO = "anthropic/claude-3.5-sonnet"
    
    # Modelos alternativos (descomenta el que prefieras):
    # MODELO_TEXTO = "openai/gpt-4-turbo"              # Muy buena calidad
    # MODELO_TEXTO = "openai/gpt-3.5-turbo"            # Rápido y económico
    # MODELO_TEXTO = "anthropic/claude-3-haiku"        # Muy rápido y barato
    # MODELO_TEXTO = "meta-llama/llama-3.1-70b-instruct"  # GRATIS

    prompt_completo = f"""
    Eres un narrador experto que convierte contenido educativo en cuentos fascinantes.
    
    Tu tarea:
    1. Lee y comprende el siguiente contenido.
    2. Crea una historia o cuento narrativo que explique los conceptos principales de manera entretenida.
    3. La narración debe ser como un cuento, con introducción, desarrollo y conclusión.
    4. Usa metáforas, personajes o situaciones que hagan el contenido memorable.
    5. Mantén un tono amigable y educativo.
    6. La narración debe durar aproximadamente 2-3 minutos al ser leída (400-600 palabras).
    
    Contenido a transformar en cuento:
    {contenido_archivo[:4000]}  # Limitamos el contenido para no exceder tokens
    
    Genera solo el cuento narrativo, sin introducciones adicionales.
    """

    try:
        respuesta_chat = client_text.chat.completions.create(
            model=MODELO_TEXTO,
            messages=[
                {"role": "system", "content": "Eres un narrador experto que transforma contenido educativo en cuentos fascinantes y memorables."},
                {"role": "user", "content": prompt_completo}
            ],
            max_tokens=1500,
            temperature=0.7
        )
        texto_para_audio = respuesta_chat.choices[0].message.content
        print("✅ Texto generado con éxito por OpenRouter.")
        
    except Exception as e:
        print(f"❌ Error en la llamada a OpenRouter: {e}")
        raise Exception(f"Error generando texto: {e}")

    # 2. --- ETAPA DE AUDIO: Conversión con OpenAI (TTS-1) ---
    print("Convirtiendo texto generado a MP3 con TTS-1...")
    speech_file_path = Path(app.config['OUTPUT_FOLDER']) / nombre_salida 
    
    try:
        with client_audio.audio.speech.with_streaming_response.create(
            model="tts-1-hd",
            voice="nova",
            input=texto_para_audio,
        ) as response:
            response.stream_to_file(speech_file_path)
        
        print(f"✅ Proceso completado. Archivo MP3 guardado en: {speech_file_path}")
        return str(speech_file_path), texto_para_audio
        
    except Exception as e:
        print(f"❌ Error en la llamada a TTS-1: {e}")
        raise Exception(f"Error generando audio: {e}")

# --- RUTAS DE LA APLICACIÓN WEB ---

@app.route('/')
def index():
    """Página principal con la interfaz de carga."""
    return render_template('index.html')

@app.route('/procesar', methods=['POST'])
def procesar():
    """Procesa el archivo o texto y genera el MP3."""
    try:
        contenido = ""
        
        # Opción 1: Texto directo desde el chat
        if 'texto_directo' in request.form and request.form['texto_directo'].strip():
            contenido = request.form['texto_directo']
            print("📝 Procesando texto directo...")
        
        # Opción 2: Archivo subido
        elif 'archivo' in request.files:
            file = request.files['archivo']
            if file and file.filename != '' and allowed_file(file.filename):
                # Guardar archivo temporalmente
                filename = secure_filename(file.filename)
                unique_filename = f"{uuid.uuid4()}_{filename}"
                filepath = os.path.join(app.config['UPLOAD_FOLDER'], unique_filename)
                file.save(filepath)
                
                print(f"📁 Procesando archivo: {filename}")
                contenido = extraer_texto(filepath)
                
                # Eliminar archivo temporal
                os.remove(filepath)
            else:
                return jsonify({'error': 'Archivo no válido. Formatos permitidos: PDF, DOCX, PPTX, TXT'}), 400
        else:
            return jsonify({'error': 'Debes proporcionar un archivo o texto directo'}), 400
        
        # Validar que se extrajo contenido
        if not contenido or contenido.startswith("Error") or len(contenido) < 50:
            return jsonify({'error': 'No se pudo extraer suficiente contenido del archivo o texto'}), 400
        
        # Generar nombre descriptivo para el audio con fecha y hora
        from datetime import datetime
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        audio_filename = f"cuento_{timestamp}.mp3"
        
        # Procesar y generar audio
        audio_path, texto_generado = procesar_y_generar_audio(contenido, audio_filename)
        
        return jsonify({
            'success': True,
            'audio_url': f'/descargar/{audio_filename}',
            'texto_generado': texto_generado,
            'filename': audio_filename
        })
        
    except Exception as e:
        print(f"❌ Error en procesamiento: {e}")
        return jsonify({'error': str(e)}), 500

@app.route('/descargar/<filename>')
def descargar(filename):
    """Descarga el archivo MP3 generado."""
    try:
        file_path = os.path.join(app.config['OUTPUT_FOLDER'], filename)
        return send_file(file_path, as_attachment=True)
    except Exception as e:
        return jsonify({'error': 'Archivo no encontrado'}), 404

# --- EJECUCIÓN ---
if __name__ == "__main__":
    # Modo de desarrollo con auto-reload
    app.run(debug=True, host='0.0.0.0', port=5000)